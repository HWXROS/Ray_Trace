"""
生成课程设计报告 (Word) 和汇报 PPT
"""
import sys
sys.path.insert(0, '.')

from docx import Document
from docx.shared import Inches, Pt, RGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
from docx.enum.style import WD_STYLE_TYPE
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.util import Inches as PptxInches
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

OUTPUT_DIR = 'outputs'
REPORT_PATH = 'report.docx'
PPT_PATH = 'report.pptx'

# ============================================================
# Word 报告生成
# ============================================================

def set_chinese_font(run, font_name='宋体', size=12, bold=False):
    font = run.font
    font.name = font_name
    font.size = Pt(size)
    font.bold = bold
    run._element.rPr.rFonts.set(qn('w:eastAsia'), font_name)

def add_heading_zh(doc, text, level=1):
    """添加中文标题"""
    # 使用英文样式名避免编码问题
    style_map = {1: 'Heading 1', 2: 'Heading 2', 3: 'Heading 3'}
    style_name = style_map.get(level, 'Heading 3')
    p = doc.add_paragraph(style=style_name)
    run = p.add_run(text)
    set_chinese_font(run, '黑体', 16 if level == 1 else 14 if level == 2 else 12, bold=True)
    return p

def add_paragraph_zh(doc, text, indent=True, font_name='宋体', size=12, bold=False, alignment=WD_ALIGN_PARAGRAPH.LEFT):
    """添加中文正文段落"""
    p = doc.add_paragraph()
    if indent:
        p.paragraph_format.first_line_indent = Cm(0.74)
    p.paragraph_format.line_spacing_rule = WD_LINE_SPACING.ONE_POINT_FIVE
    p.alignment = alignment
    run = p.add_run(text)
    set_chinese_font(run, font_name, size, bold)
    return p

def add_image_zh(doc, img_path, width_cm=14, caption=None):
    """添加图片及说明"""
    if os.path.exists(img_path):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run()
        run.add_picture(img_path, width=Cm(width_cm))
        if caption:
            cap_p = doc.add_paragraph()
            cap_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            cap_run = cap_p.add_run(f'图 {caption}')
            set_chinese_font(cap_run, '宋体', 10)
        return True
    return False

def generate_report():
    doc = Document()
    
    # 设置默认中文字体
    style = doc.styles['Normal']
    style.font.name = '宋体'
    style._element.rPr.rFonts.set(qn('w:eastAsia'), '宋体')
    style.font.size = Pt(12)
    
    # ========== 封面 ==========
    for _ in range(6):
        doc.add_paragraph()
    
    title_p = doc.add_paragraph()
    title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_run = title_p.add_run('高级图形学与增强现实')
    set_chinese_font(title_run, '黑体', 22, bold=True)
    
    title_p2 = doc.add_paragraph()
    title_p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_run2 = title_p2.add_run('光线追踪器课程设计报告')
    set_chinese_font(title_run2, '黑体', 22, bold=True)
    
    for _ in range(4):
        doc.add_paragraph()
    
    info_lines = [
        '西安交通大学 人工智能学院',
        '姓    名：XXX',
        '学    号：XXXXXXXXXX',
        '完成日期：2026年5月',
    ]
    for line in info_lines:
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(line)
        set_chinese_font(run, '宋体', 14)
    
    doc.add_page_break()
    
    # ========== 目录（占位） ==========
    add_heading_zh(doc, '目录', level=1)
    toc_items = [
        '1  内容与要求',
        '    1.1  选题背景',
        '    1.2  功能需求',
        '2  总体方案',
        '    2.1  系统架构',
        '    2.2  核心算法',
        '    2.3  技术路线',
        '3  实验结果',
        '    3.1  基础场景渲染',
        '    3.2  材质与图元展示',
        '    3.3  特色场景',
        '    3.4  消融实验',
        '    3.5  性能分析',
        '4  心得体会',
    ]
    for item in toc_items:
        p = doc.add_paragraph()
        run = p.add_run(item)
        set_chinese_font(run, '宋体', 12)
    doc.add_page_break()
    
    # ========== 第1章 内容与要求 ==========
    add_heading_zh(doc, '1  内容与要求', level=1)
    
    add_heading_zh(doc, '1.1  选题背景', level=2)
    add_paragraph_zh(doc, 
        '本课程设计选择题目为《光线追踪器》，参考 Peter Shirley 的《Ray Tracing in One Weekend》系列书籍，'
        '基于 Python + Taichi 框架实现了一套完整的 GPU 加速光线追踪渲染器。'
        '光线追踪（Ray Tracing）是计算机图形学中生成照片级真实感图像的核心算法，'
        '通过模拟光线在场景中的传播路径（反射、折射、散射）来计算像素颜色。'
        '相比传统的光栅化渲染，光线追踪能够自然地呈现全局光照、软阴影、反射折射等复杂光学现象。')
    
    add_heading_zh(doc, '1.2  功能需求', level=2)
    add_paragraph_zh(doc, '根据课程设计要求，将光线追踪器的功能分解为以下子模块：')
    
    features = [
        ('核心光追算法', [
            '光线与图元相交：支持球体、立方体（AABB）、三角形三种基本图元',
            'Lambertian 漫反射材质：基于 Cosine-weighted 半球采样',
            'Metal 金属材质：镜面反射 + 可调模糊度（fuzz）',
            'Dielectric 电介质材质：Fresnel 折射（Schlick 近似），支持不同折射率',
            'Perlin 噪声纹理：程序化生成大理石等自然纹理',
            'Emissive 发光材质：区域光源，支持 NEE 直接采样',
            '路径追踪（Path Tracing）：Monte Carlo 积分 + Russian Roulette 截断',
            'NEE（Next Event Estimation）：直接光源采样加速收敛',
        ]),
        ('加速与优化', [
            'BVH 加速结构：CPU 端构建，GPU 端 skip-link 无栈遍历',
            'Taichi CUDA Kernel：每个像素独立 GPU 线程并行渲染',
            'Gamma 校正：物理正确到显示正确的颜色转换',
        ]),
        ('高级特性', [
            'OBJ 模型加载：通过 trimesh 加载外部模型并接入渲染管线',
            'HDRI 环境贴图：equirectangular 球形映射实现真实环境光照',
            'YAML 场景配置：无需修改代码即可定义场景、相机和动画',
            '动画渲染：相机路径关键帧插值生成帧序列',
            '交互式预览：ti.GUI 实时显示渐进渲染过程',
        ]),
        ('真实感效果', [
            '软阴影：区域光源自然产生的半影效果',
            '景深（Defocus Blur）：薄透镜模型模拟相机光圈虚化',
            '抗锯齿：每像素多采样 + 随机抖动',
        ]),
    ]
    
    for title, items in features:
        p = doc.add_paragraph()
        p.paragraph_format.first_line_indent = Cm(0.74)
        run = p.add_run(f'（1）{title}：')
        set_chinese_font(run, '宋体', 12, bold=True)
        for item in items:
            p = doc.add_paragraph()
            p.paragraph_format.first_line_indent = Cm(1.27)
            run = p.add_run(f'• {item}')
            set_chinese_font(run, '宋体', 12)
    
    doc.add_page_break()
    
    # ========== 第2章 总体方案 ==========
    add_heading_zh(doc, '2  总体方案', level=1)
    
    add_heading_zh(doc, '2.1  系统架构', level=2)
    add_paragraph_zh(doc,
        '系统采用分层架构设计，分为 GPU 渲染核心层、场景管理层和工具脚本层。'
        'GPU 核心层使用 Taichi 的 @ti.kernel 和 @ti.func 实现 CUDA 并行渲染；'
        '场景管理层提供 Python API 构建场景并配置相机；'
        '工具脚本层提供命令行入口、YAML 解析和批量渲染功能。')
    add_paragraph_zh(doc,
        '项目目录结构如下：renderer/ 存放 GPU 核心代码（taichi_renderer.py、perlin.py、config.py）；'
        'scenes/ 存放场景预设和 YAML 配置；scripts/ 存放运行入口；'
        'experiments/ 存放消融实验和基准测试；geometry/ 提供 OBJ 加载和三角形支持。')
    
    add_heading_zh(doc, '2.2  核心算法', level=2)
    add_paragraph_zh(doc,
        '（1）路径追踪循环：每条光线从相机发出，在场景中递归（循环模拟）追踪，'
        '每次命中表面后根据材质类型采样散射方向。深度超过 3 后启用 Russian Roulette，'
        '按生存概率终止低贡献路径。')
    add_paragraph_zh(doc,
        '（2）NEE 直接光源采样：对 Lambertian 材质，在前 3 次反弹时向随机光源球面采样一点，'
        '发送 Shadow Ray 测试可见性，将直接光照贡献加入最终颜色。PDF 与 BRDF 权重保证无偏性。')
    add_paragraph_zh(doc,
        '（3）BVH 加速遍历：CPU 端递归构建 AABB 层次包围盒，DFS 前序线性化存储为 Taichi field。'
        'GPU 端使用 skip-link 无栈遍历：每个节点存储 skip 指针指向子树结束后的下一个节点，'
        'AABB 不相交时直接跳过整个子树。')
    add_paragraph_zh(doc,
        '（4）Fresnel 折射：电介质材质根据 Schlick 近似计算反射概率 R(θ) = R₀ + (1-R₀)(1-cosθ)⁵，'
        '按概率选择反射或折射路径，模拟真实玻璃的行为。')
    
    add_heading_zh(doc, '2.3  技术路线', level=2)
    add_paragraph_zh(doc,
        '开发分为四个阶段：第一阶段实现基础光线追踪（球体相交、Lambertian/Metal 材质、简单场景）；'
        '第二阶段引入 GPU 加速（Taichi CUDA、BVH 构建与遍历）；'
        '第三阶段增加高级材质和光照（Dielectric、Emissive、NEE、Russian Roulette）；'
        '第四阶段完善工程能力（OBJ 加载、HDRI、YAML 驱动、动画、交互预览）。')
    
    doc.add_page_break()
    
    # ========== 第3章 实验结果 ==========
    add_heading_zh(doc, '3  实验结果', level=1)
    
    add_heading_zh(doc, '3.1  基础场景渲染', level=2)
    add_paragraph_zh(doc,
        '基础三球场景（taichi_demo.png）验证了 GPU 渲染管线的正确性，包含红色 Lambertian 球、'
        '银色金属球和玻璃球（Fresnel 折射）。复杂随机场景（taichi_complex.png）包含 486 个随机球体，'
        '展示了大规模场景的渲染能力和大光圈景深效果。')
    add_image_zh(doc, f'{OUTPUT_DIR}/taichi_demo.png', width_cm=12, caption='基础三球场景')
    add_image_zh(doc, f'{OUTPUT_DIR}/taichi_complex.png', width_cm=14, caption='复杂随机场景（486球+景深）')
    
    add_heading_zh(doc, '3.2  材质与图元展示', level=2)
    add_paragraph_zh(doc,
        '材质博览馆场景（showcase_01）将 5 种材质和 3 种图元按矩阵排列，'
        '清晰展示了 Lambertian（漫反射）、Metal（金属，不同 fuzz）、Dielectric（玻璃）、'
        'Perlin（噪声纹理）、Emissive（发光）的视觉效果差异。')
    add_image_zh(doc, f'{OUTPUT_DIR}/showcase_01_materials.png', width_cm=12, caption='材质博览馆')
    
    add_paragraph_zh(doc,
        '金属工坊（showcase_02）聚焦于不同金属颜色（银、金、铜）和模糊度参数的反射表现；'
        '水晶宫（showcase_03）展示了多种折射率（玻璃 1.5、蓝宝石 1.77、钻石 2.42、水 1.33）的折射差异，'
        '以及玻璃立方体、球体和金字塔的形态。')
    add_image_zh(doc, f'{OUTPUT_DIR}/showcase_02_metals.png', width_cm=12, caption='金属工坊')
    add_image_zh(doc, f'{OUTPUT_DIR}/showcase_03_glass.png', width_cm=12, caption='水晶宫（玻璃折射）')
    
    add_paragraph_zh(doc,
        '霓虹展台（showcase_04）使用 Emissive 发光立方体和三角形作为区域光源，'
        '在被照亮的金属球和漫反射表面形成柔和的彩色阴影。')
    add_image_zh(doc, f'{OUTPUT_DIR}/showcase_04_neon.png', width_cm=12, caption='霓虹展台（发光材质）')
    
    add_heading_zh(doc, '3.3  特色场景', level=2)
    add_paragraph_zh(doc,
        '几何森林（showcase_05）在绿色地面上随机散布 30 个混合图元，'
        '展示渲染器处理大量不同几何体的能力；'
        'Cornell Box（showcase_06）是经典的封闭空间全局光照测试场景，'
        '红绿墙壁的间接光照颜色渗透在白球表面清晰可见；'
        '创新港微缩模型（showcase_07）参照西安交通大学创新港校区平面图，'
        '用 86 个立方体和 41 个球体搭建简化版校园沙盘，'
        '黄色为教学楼、粉色为宿舍、绿色为绿地，从 45° 鸟瞰视角呈现中轴线对称布局。')
    add_image_zh(doc, f'{OUTPUT_DIR}/showcase_05_forest.png', width_cm=12, caption='几何森林')
    add_image_zh(doc, f'{OUTPUT_DIR}/showcase_06_cornell.png', width_cm=12, caption='Cornell Box（全局光照）')
    add_image_zh(doc, f'{OUTPUT_DIR}/showcase_07_innovation_harbor.png', width_cm=14, caption='创新港微缩模型')
    
    add_heading_zh(doc, '3.4  消融实验', level=2)
    add_paragraph_zh(doc,
        '（1）分辨率/采样/光圈消融实验（ablation_study_v2.png）：'
        '第一行固定光圈和采样，展示分辨率从 200×200 到 800×800 的清晰度变化；'
        '第二行固定分辨率和光圈，展示采样数从 10spp 到 500spp 对 Monte Carlo 噪点的影响，'
        '10spp 时噪点严重，500spp 图像基本纯净；'
        '第三行固定分辨率和采样，展示光圈从 0 到 2.0 的景深效果，'
        '光圈越大，失焦区域越模糊，中景金属球保持清晰。')
    add_image_zh(doc, f'{OUTPUT_DIR}/ablation_study_v2.png', width_cm=14, caption='分辨率/采样/光圈消融实验')
    
    add_paragraph_zh(doc,
        '（2）封闭空间消融实验（ablation_enclosed_v2.png）：'
        '第一行展示采样数对封闭 Cornell Box 收敛速度的影响，50spp 噪点明显，1000spp 接近收敛；'
        '第二行展示光源面积对软阴影的影响，光源越大阴影越柔和；'
        '第三行展示封闭程度对光照的影响，从 1 面（开放）到 6 面（全封闭），'
        '全封闭时仅靠内部球光源照明，颜色渗透和间接光照效果最显著。')
    add_image_zh(doc, f'{OUTPUT_DIR}/ablation_enclosed_v2.png', width_cm=14, caption='封闭空间消融实验')
    
    add_paragraph_zh(doc,
        '（3）BVH/NEE 性能消融实验（ablation_bvh_nee.png）：'
        '在小（~15 图元）、中（~80 图元）、大（~300 图元）三个场景下，'
        '分别测试 BVH+NEE、BVH only、NEE only、None 四种组合的渲染时间。'
        '实验发现当前 skip-link BVH 在 GPU 上因 warp divergence 导致性能不如线性遍历，'
        '这是 GPU BVH 遍历的已知挑战；NEE 增加了 shadow ray 开销，但在封闭空间中能显著降低噪点。')
    add_image_zh(doc, f'{OUTPUT_DIR}/ablation_bvh_nee.png', width_cm=14, caption='BVH/NEE 性能消融实验')
    
    add_heading_zh(doc, '3.5  性能分析', level=2)
    add_paragraph_zh(doc,
        '表 1 列出了 CPU 版本与 GPU 版本的性能对比。CPU 版本使用纯 Python 实现，'
        '渲染 400×225、100spp 的 Demo 场景需要约 127 秒；'
        'GPU 版本（Taichi CUDA）仅需约 0.95 秒，加速比约 150 倍。'
        '在 800×450、500spp 的高质量设置下，Complex 场景渲染约 16 秒。')
    
    # 性能表格
    table = doc.add_table(rows=1, cols=5)
    table.style = 'Light Grid Accent 1'
    hdr_cells = table.rows[0].cells
    headers = ['版本', '场景', '分辨率', '采样', '耗时']
    for i, h in enumerate(headers):
        hdr_cells[i].text = h
        for p in hdr_cells[i].paragraphs:
            for r in p.runs:
                set_chinese_font(r, '黑体', 11, bold=True)
    
    rows_data = [
        ['CPU (Python)', 'Demo', '400×225', '100', '~127s'],
        ['GPU (Taichi)', 'Demo', '400×225', '100', '~0.95s'],
        ['GPU (Taichi)', 'Complex', '800×450', '100', '~5.6s'],
        ['GPU (Taichi)', 'Complex', '800×450', '500', '~16s'],
    ]
    for row_data in rows_data:
        row_cells = table.add_row().cells
        for i, val in enumerate(row_data):
            row_cells[i].text = val
            for p in row_cells[i].paragraphs:
                for r in p.runs:
                    set_chinese_font(r, '宋体', 11)
    
    add_paragraph_zh(doc, '表 1  CPU vs GPU 渲染性能对比', indent=False)
    add_paragraph_zh(doc,
        'BVH 构建时间上，21 节点的小场景构建耗时 <1ms，503 节点的大场景构建耗时约 5ms，'
        'CPU 端构建开销可忽略不计。渲染过程中发现的 build_bvh() bug（每次创建新 field 导致 kernel 重新编译）'
        '已修复为复用已有 field，消除了约 17 秒的额外编译开销。')
    
    doc.add_page_break()
    
    # ========== 第4章 心得体会 ==========
    add_heading_zh(doc, '4  心得体会', level=1)
    add_paragraph_zh(doc,
        '通过本次课程设计，深入理解了光线追踪算法的原理与实现细节。'
        '路径追踪虽然概念简单——"从相机发出光线，在场景中反复弹跳"——'
        '但实际实现中涉及大量工程细节：蒙特卡洛积分的无偏性保证、'
        'Russian Roulette 的生存概率选择、NEE 的 PDF 权重计算、'
        'Fresnel 折射的 Schlick 近似等。')
    add_paragraph_zh(doc,
        'GPU 并行化是最大的挑战。Taichi 框架极大地简化了 CUDA Kernel 的编写，'
        '但性能调优仍然需要深入理解 GPU 执行模型。'
        'BVH 的 skip-link 无栈遍历在理论上很优雅，但在 GPU 上因 warp divergence 导致实际性能不佳，'
        '这让我们认识到"算法复杂度低"不等于"GPU 上跑得快"，'
        '内存访问模式和线程同步策略同样重要。')
    add_paragraph_zh(doc,
        '此外，工程能力同样关键。从 YAML 配置驱动到交互式预览，'
        '从 OBJ 模型加载到动画渲染，这些"外围功能"让算法从 demo 变成了可用的工具。'
        '特别是消融实验的设计，让我们学会了用控制变量的方法验证每个技术点的实际效果。')
    add_paragraph_zh(doc,
        '最后，创新港微缩场景的搭建让我们体会到图形学不仅是算法，'
        '也是艺术表达的工具。用代码"建造"一座虚拟校园，'
        '看到沙盘模型在 GPU 上逐渐清晰，是一种独特的成就感。')
    
    doc.save(REPORT_PATH)
    print(f'[Report] Word 报告已生成: {REPORT_PATH}')
    return doc


# ============================================================
# PPT 生成
# ============================================================

def generate_ppt():
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)
    
    def add_title_slide(title, subtitle=''):
        layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(layout)
        
        title_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(2), PptxInches(12.3), PptxInches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.name = '微软雅黑'
        p.alignment = PP_ALIGN.CENTER
        
        if subtitle:
            sub_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(3.8), PptxInches(12.3), PptxInches(1))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(20)
            p.font.name = '微软雅黑'
            p.alignment = PP_ALIGN.CENTER
        return slide
    
    def add_content_slide(title, bullets, img_path=None, img_left=None, img_top=None, img_width=None):
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)
        
        title_box = slide.shapes.add_textbox(PptxInches(0.4), PptxInches(0.3), PptxInches(12.5), PptxInches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.name = '微软雅黑'
        
        bullet_box = slide.shapes.add_textbox(PptxInches(0.4), PptxInches(1.2), PptxInches(12.5), PptxInches(5.5))
        tf = bullet_box.text_frame
        tf.word_wrap = True
        for i, text in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = text
            p.font.size = Pt(20)
            p.font.name = '微软雅黑'
            p.space_after = Pt(12)
        
        if img_path and os.path.exists(img_path) and img_left is not None:
            slide.shapes.add_picture(img_path, img_left, img_top, width=img_width)
        
        return slide
    
    def add_image_slide(title, img_paths, captions=None):
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)
        
        title_box = slide.shapes.add_textbox(PptxInches(0.4), PptxInches(0.2), PptxInches(12.5), PptxInches(0.7))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.name = '微软雅黑'
        
        n = len(img_paths)
        if n == 1:
            positions = [(PptxInches(3.5), PptxInches(1.1), PptxInches(6.3))]
        elif n == 2:
            positions = [(PptxInches(0.5), PptxInches(1.1), PptxInches(6)),
                        (PptxInches(6.8), PptxInches(1.1), PptxInches(6))]
        elif n == 3:
            positions = [(PptxInches(0.3), PptxInches(1.1), PptxInches(4)),
                        (PptxInches(4.6), PptxInches(1.1), PptxInches(4)),
                        (PptxInches(9.0), PptxInches(1.1), PptxInches(4))]
        else:
            positions = [(PptxInches(0.3), PptxInches(1.1), PptxInches(3)),
                        (PptxInches(3.5), PptxInches(1.1), PptxInches(3)),
                        (PptxInches(6.7), PptxInches(1.1), PptxInches(3)),
                        (PptxInches(9.9), PptxInches(1.1), PptxInches(3))]
        
        for i, (img_path, (left, top, width)) in enumerate(zip(img_paths, positions)):
            if os.path.exists(img_path):
                slide.shapes.add_picture(img_path, left, top, width=width)
                if captions and i < len(captions):
                    cap_box = slide.shapes.add_textbox(left, top + width * 0.6 + PptxInches(0.1), width, PptxInches(0.4))
                    tf = cap_box.text_frame
                    p = tf.paragraphs[0]
                    p.text = captions[i]
                    p.font.size = Pt(12)
                    p.font.name = '微软雅黑'
                    p.alignment = PP_ALIGN.CENTER
        return slide
    
    # 1. 封面
    add_title_slide('GPU 光线追踪器', '高级图形学与增强现实 课程设计报告')
    
    # 2. 项目概述
    add_content_slide('项目概述', [
        '选题：光线追踪器（参考《Ray Tracing in One Weekend》）',
        '技术栈：Python + Taichi + CUDA',
        'GPU：NVIDIA RTX 4060Ti 8G',
        '核心特性：路径追踪、BVH 加速、NEE 直接采样、多种材质',
        '加速比：CPU vs GPU ≈ 150 倍',
        '扩展能力：OBJ 加载、HDRI、YAML 配置、动画渲染',
    ])
    
    # 3. 系统架构
    add_content_slide('系统架构', [
        'renderer/ — GPU 渲染核心（Taichi Kernel、Perlin 噪声、YAML 解析）',
        'scenes/ — 场景预设（7 个展示场景 + YAML 配置）',
        'scripts/ — 运行入口（命令行、拼图、动画）',
        'experiments/ — 消融实验（分辨率/采样/光圈/封闭空间/BVH-NEE）',
        'geometry/ — OBJ 模型加载 + 程序化立方体',
    ])
    
    # 4. 核心算法
    add_content_slide('核心算法', [
        '路径追踪：Monte Carlo 积分 + Russian Roulette 截断',
        'NEE：向光源球面采样 + Shadow Ray 可见性测试',
        'BVH：CPU 构建 + DFS 线性化 + GPU skip-link 无栈遍历',
        'Fresnel 折射：Schlick 近似 R(θ) = R₀ + (1-R₀)(1-cosθ)⁵',
        'Perlin 噪声：程序化生成大理石等自然纹理',
    ])
    
    # 5. 材质展示
    add_image_slide('材质与图元展示',
        [f'{OUTPUT_DIR}/showcase_01_materials.png', f'{OUTPUT_DIR}/showcase_02_metals.png', f'{OUTPUT_DIR}/showcase_03_glass.png'],
        ['5 种材质 × 3 种图元', '金属反射（不同 fuzz）', '玻璃折射（不同折射率）'])
    
    # 6. 特色场景
    add_image_slide('特色场景',
        [f'{OUTPUT_DIR}/showcase_04_neon.png', f'{OUTPUT_DIR}/showcase_05_forest.png', f'{OUTPUT_DIR}/showcase_06_cornell.png'],
        ['霓虹展台（发光材质）', '几何森林（大量图元）', 'Cornell Box（全局光照）'])
    
    # 7. 创新港场景
    add_content_slide('创新港微缩模型', [
        '参照西安交通大学创新港校区平面图搭建',
        '扇形倒三角布局，北部宽南部窄',
        '黄色=教学楼（泓理楼/躬行楼/力行楼/涵英楼/敏行楼）',
        '粉色=学生宿舍（B区/C区/A区/Y形宿舍）',
        '绿色=绿地，深绿=运动场地',
        '86 个立方体 + 41 个球体，45° 沙盘视角',
    ], f'{OUTPUT_DIR}/showcase_07_innovation_harbor.png', PptxInches(7.5), PptxInches(0.8), PptxInches(5.5))
    
    # 8. 消融实验1
    add_image_slide('消融实验 — 分辨率/采样/光圈',
        [f'{OUTPUT_DIR}/ablation_study_v2.png'],
        ['分辨率↑清晰度↑ | 采样↑噪点↓ | 光圈↑景深虚化↑'])
    
    # 9. 消融实验2
    add_image_slide('消融实验 — 封闭空间 & BVH/NEE',
        [f'{OUTPUT_DIR}/ablation_enclosed_v2.png', f'{OUTPUT_DIR}/ablation_bvh_nee.png'],
        ['采样/光源/封闭程度影响', 'BVH/NEE 性能分析'])
    
    # 10. 性能数据
    add_content_slide('性能分析', [
        'CPU (Python)  Demo 400×225 100spp：~127s',
        'GPU (Taichi)  Demo 400×225 100spp：~0.95s',
        '加速比：≈ 150 倍',
        'Complex 800×450 500spp：~16s',
        '修复 build_bvh field 替换 bug，消除 ~17s 编译开销',
    ])
    
    # 11. 总结
    add_content_slide('总结与展望', [
        '实现了完整的 GPU 加速路径追踪渲染器',
        '支持 5 种材质、3 种图元、BVH 加速、NEE、HDRI、OBJ 加载',
        '消融实验验证了各技术参数对画质和性能的影响',
        '未来优化方向：',
        '  • GPU BVH 遍历改用栈式或重新排序光线以减少 warp divergence',
        '  • 引入多重重要性采样（MIS）进一步提升 NEE 效率',
        '  • 支持纹理贴图和更复杂的 BSDF 模型',
    ])
    
    prs.save(PPT_PATH)
    print(f'[PPT] 汇报 PPT 已生成: {PPT_PATH}')
    return prs


if __name__ == '__main__':
    generate_report()
    generate_ppt()
    print('\n全部完成！')
